"""PowerPoint generation — professional decks with charts, images, timelines, and 10+ layouts."""

from __future__ import annotations

import json
import logging
import os
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any

import httpx
from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.agents.tools.registry import ToolContext, ToolHandler
from src.types.agent import ToolDefinition

logger = logging.getLogger("kinetic.presentation")

SANDBOX = Path("agent_sandbox")


@dataclass
class StylePreset:
    name: str
    accent: RGBColor
    dark: RGBColor
    light: RGBColor
    bg: RGBColor
    title_size: int = 44
    slide_title: int = 28
    body: int = 18
    code: int = 13


STYLES: dict[str, StylePreset] = {
    "corporate": StylePreset(
        "Corporate",
        RGBColor(0x1A, 0x73, 0xE8),
        RGBColor(0x20, 0x20, 0x20),
        RGBColor(0xCC, 0xDD, 0xFF),
        RGBColor(0xFF, 0xFF, 0xFF),
    ),
    "dark-tech": StylePreset(
        "Dark Tech",
        RGBColor(0x00, 0xCC, 0xFF),
        RGBColor(0xE0, 0xE0, 0xE0),
        RGBColor(0x00, 0x66, 0x99),
        RGBColor(0x1A, 0x1A, 0x2E),
    ),
    "minimal": StylePreset(
        "Minimal",
        RGBColor(0x33, 0x33, 0x33),
        RGBColor(0x11, 0x11, 0x11),
        RGBColor(0x99, 0x99, 0x99),
        RGBColor(0xFA, 0xFA, 0xFA),
    ),
    "warm": StylePreset(
        "Warm",
        RGBColor(0xE8, 0x6C, 0x00),
        RGBColor(0x33, 0x22, 0x11),
        RGBColor(0xFF, 0xDD, 0xAA),
        RGBColor(0xFF, 0xFD, 0xF5),
    ),
    "nature": StylePreset(
        "Nature",
        RGBColor(0x2E, 0x7D, 0x32),
        RGBColor(0x1B, 0x33, 0x1B),
        RGBColor(0xA5, 0xD6, 0xA7),
        RGBColor(0xFB, 0xFF, 0xF9),
    ),
    "ocean": StylePreset(
        "Ocean",
        RGBColor(0x00, 0x78, 0x9C),
        RGBColor(0x02, 0x23, 0x33),
        RGBColor(0x9E, 0xCF, 0xE2),
        RGBColor(0xF0, 0xF8, 0xFF),
    ),
    "sunset": StylePreset(
        "Sunset",
        RGBColor(0x9C, 0x27, 0xB0),
        RGBColor(0x2C, 0x0B, 0x33),
        RGBColor(0xCE, 0x93, 0xD8),
        RGBColor(0xFF, 0xF3, 0xE0),
    ),
    "forest": StylePreset(
        "Forest",
        RGBColor(0x55, 0x8B, 0x2F),
        RGBColor(0x1B, 0x2E, 0x0E),
        RGBColor(0xA0, 0xD4, 0x6F),
        RGBColor(0xF5, 0xFB, 0xEE),
    ),
}

CANVAS_FORMATS: dict[str, tuple[Any, Any]] = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
    "a4": (Inches(11.69), Inches(8.27)),
}

TRANSITIONS: dict[str, str] = {
    "fade": "fade",
    "push": "push",
    "wipe": "wipe",
    "cover": "cover",
    "uncover": "uncover",
    "dissolve": "dissolve",
}

# ── Helpers ──


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def _tb(slide, left: float, t: float, w: float, h: float) -> Any:
    return slide.shapes.add_textbox(Inches(left), Inches(t), Inches(w), Inches(h))


def _accent_bar(slide, w: int, st: StylePreset) -> Any:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), w, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = st.accent
    bar.line.fill.background()
    return bar


def _title_tf(bar: Any, text: str, st: StylePreset) -> None:
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(st.slide_title)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if st.name != "Minimal" else st.dark
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)


def _bullets(tf, items: list[str], st: StylePreset, body_size: int | None = None) -> None:
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if item.startswith("#"):
            p.text = item.lstrip("#").strip()
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = st.accent
        elif item.startswith(">"):
            p.text = item.lstrip(">").strip()
            p.font.size = Pt(14)
            p.font.color.rgb = st.light
            p.font.italic = True
        else:
            p.text = f"  {item}"
            p.font.size = Pt(body_size or st.body)
            p.font.color.rgb = st.dark
        p.space_after = Pt(6)


def _detect_layout(title: str, content: list[str]) -> str:
    text = " ".join(content).lower()
    if any(w in text for w in ["question", "quiz", "?"]) and any(c.startswith("?") for c in content):
        return "quiz"
    if "```" in text or any(kw in text for kw in ["def ", "import ", "class ", "return "]):
        return "code"
    if len(content) > 6 or any(len(c) > 100 for c in content):
        return "two_column"
    if len(content) == 1 and len(content[0]) < 60:
        return "center_text"
    return "bullets"


async def _dl_img(src: str) -> bytes | None:
    if src.startswith(("http://", "https://")):
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                r = await client.get(src, follow_redirects=True)
                r.raise_for_status()
                return r.content
        except Exception:
            return None
    for base in [SANDBOX, Path(".")]:
        p = (base / src).resolve()
        if p.exists() and p.is_file():
            return p.read_bytes()
    return None


def _place_img(slide, data: bytes, pos: str, sw: int, sh: int):
    try:
        img = slide.shapes.add_picture(BytesIO(data), Inches(0), Inches(0))
    except Exception:
        return None, None
    iw, ih = img.width / 914400, img.height / 914400
    swi, shi = sw / 914400, sh / 914400

    if pos == "bg":
        img.width = sw
        img.height = sh
        img.left = 0
        img.top = 0
        # send to back
        sp = img._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
        return None, None

    max_w = swi * 0.45 if pos != "full" else swi - 1.0
    max_h = shi - 1.2 if pos != "full" else shi - 1.5
    scale = min(max_w / iw, max_h / ih, 1.0)
    fw, fh = int(iw * scale * 914400), int(ih * scale * 914400)
    if pos == "right":
        img.left = int((swi - swi * 0.48) * 914400)
        img.top = int(1.1 * 914400)
    elif pos == "left":
        img.left = int(0.3 * 914400)
        img.top = int(1.1 * 914400)
    else:  # full
        img.left = int((swi - fw / 914400) / 2 * 914400)
        img.top = int(1.2 * 914400)
    img.width, img.height = fw, fh
    return fw / 914400, fh / 914400


async def _build_slide(
    prs,
    title: str,
    content: list[str],
    layout: str | None = None,
    notes: str | None = None,
    st: StylePreset | None = None,
    sw: int = 0,
    sh: int = 0,
    img_src: str | None = None,
    img_pos: str = "right",
    chart_type: str | None = None,
    chart_data: list | None = None,
):
    st = st or STYLES["corporate"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if notes:
        ns = slide.notes_slide
        ns.notes_text_frame.text = notes

    detected = layout or _detect_layout(title, content)

    # ── Section divider ──
    if detected == "section":
        _set_bg(slide, st.accent)
        tf = _tb(slide, 1, 2.5, sw / 914400 - 2, 3).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        if content:
            p2 = tf.add_paragraph()
            p2.text = content[0] if isinstance(content, list) else content
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD) if st.name != "Minimal" else st.dark
            p2.alignment = PP_ALIGN.CENTER
        return

    _set_bg(slide, st.bg)
    b = _accent_bar(slide, sw, st)
    _title_tf(b, title, st)

    # Image
    img_w = None
    if img_src:
        data = await _dl_img(img_src)
        if data:
            img_w, _ = _place_img(slide, data, img_pos, sw, sh)

    y = 1.3
    m = 0.5
    bw = sw / 914400 - m * 2
    if img_src and img_pos == "right":
        bw = sw / 914400 * 0.48
    elif img_src and img_pos == "left":
        m = sw / 914400 * 0.52
        bw = sw / 914400 * 0.45

    # ── Chart ──
    if chart_type and chart_data:
        cd = CategoryChartData()
        cd.categories = [r[0] for r in chart_data]
        cd.add_series("Data", [r[1] for r in chart_data])
        xl_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        try:
            chart = slide.shapes.add_chart(
                xl_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                Inches(m),
                Inches(y),
                Inches(bw),
                Inches(5.0),
                cd,
            ).chart
            chart.has_legend = True
        except Exception:
            pass
        return

    # ── Layouts ──
    if detected == "code":
        bg = slide.shapes.add_shape(1, Inches(m - 0.2), Inches(y - 0.1), Inches(bw + 0.4), Inches(5.7))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        bg.line.fill.background()
        tf = _tb(slide, m, y, bw, 5.5).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "\n".join(content)
        p.font.size = Pt(st.code)
        p.font.name = "Consolas"
        p.font.color.rgb = st.dark

    elif detected == "comparison":
        mid = len(content) // 2
        cw = (bw - m) / 2
        for ci, (items, left) in enumerate([(content[:mid], m), (content[mid:], m + cw + m)]):
            tf = _tb(slide, left, y, cw, 5.5).text_frame
            tf.word_wrap = True
            _bullets(tf, items, st, 16)

    elif detected == "timeline":
        tf = _tb(slide, m, y, bw, 5.5).text_frame
        tf.word_wrap = True
        first = True
        for item in content:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if ":" in item and len(item.split(":")[0]) < 15:
                parts = item.split(":", 1)
                p.text = f"  {parts[0].strip()}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = st.accent
                p.space_after = Pt(2)
                p2 = tf.add_paragraph()
                p2.text = f"  {parts[1].strip()}"
                p2.font.size = Pt(16)
                p2.font.color.rgb = st.dark
                p2.space_after = Pt(10)
            elif item.startswith(">"):
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = st.light
                p.font.italic = True
            else:
                p.text = f"  {item}"
                p.font.size = Pt(st.body)
                p.font.color.rgb = st.dark
            p.space_after = Pt(6)

    elif detected == "table":
        rows = [r.split("|") for r in content if "|" in r]
        if rows and len(rows) > 1:
            ncols = max(len(r) for r in rows)
            nrows = len(rows)
            table = slide.shapes.add_table(nrows, ncols, Inches(m), Inches(y), Inches(bw), Inches(0.5 * nrows)).table
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    cell = table.cell(ri, ci)
                    cell.text = val.strip()
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        if ri == 0:
                            p.font.bold = True
                            p.font.color.rgb = st.accent
                        else:
                            p.font.color.rgb = st.dark

    elif detected == "quiz":
        tf = _tb(slide, m, y, bw, 5.5).text_frame
        tf.word_wrap = True
        first = True
        for item in content:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if item.startswith("?"):
                p.text = item
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = st.accent
            else:
                p.text = f"  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = st.dark
            p.space_after = Pt(6)

    elif detected == "center_text":
        tf = _tb(slide, m, 2.5, bw, 3).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content[0]
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = st.accent
        p.alignment = PP_ALIGN.CENTER

    elif detected == "two_column":
        mid = len(content) // 2
        cw = (bw - m) / 2
        for ci, (items, left) in enumerate([(content[:mid], m), (content[mid:], m + cw + m)]):
            tf = _tb(slide, left, y, cw, 5.5).text_frame
            tf.word_wrap = True
            first = True
            for item in items:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = f"  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = st.dark
                p.space_after = Pt(6)

    else:  # bullets
        tf = _tb(slide, m, y, bw, 5.5).text_frame
        tf.word_wrap = True
        _bullets(tf, content, st)


async def _generate_slides_from_topic(
    topic: str, title: str, model: str, api_key: str, base_url: str,
) -> list[dict[str, Any]]:
    """Use a more powerful LLM to generate slide content from a topic."""
    prompt = f"""Create a professional presentation about: {topic}

Return a JSON array of slides. Each slide has:
- "title": slide title (short, clear)
- "content": array of bullet points (3-6 each, 3-8 words each, parallel structure)

Rules:
- 4-6 slides total
- First slide after title should be an overview/agenda
- Keep content factual and specific, not vague
- Use parallel structure for bullets (same verb tense, same format)

Return ONLY valid JSON, no markdown, no explanation."""

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{base_url}/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={
                    "model": model,
                    "messages": [{"role": "user", "content": prompt}],
                    "max_completion_tokens": 2000,
                },
            )
            resp.raise_for_status()
            body = resp.text
            logger.info("[PPT] API responded: %d bytes", len(body))
            data = json.loads(body)
            raw = data["choices"][0]["message"]["content"].strip()
            logger.info("[PPT] Model response preview: %s", raw[:300])

            def _find_json(text: str) -> Any | None:
                """Extract JSON from model response using multiple strategies."""
                # 1. Direct parse
                try:
                    return json.loads(text)
                except json.JSONDecodeError:
                    pass
                # 2. Strip markdown fences
                c = text.strip()
                for p in ("```json", "```"):
                    c = c.removeprefix(p).lstrip()
                c = c.removesuffix("```").rstrip()
                try:
                    return json.loads(c)
                except json.JSONDecodeError:
                    pass
                # 3. Find balanced brackets line by line
                lines = text.split("\n")
                combined = []
                active = False
                depth = 0
                for line in lines:
                    s = line.strip()
                    if not active:
                        idx = s.find("[")
                        if idx >= 0:
                            active = True
                            combined.append(s[idx:])
                            depth = s[idx:].count("[") - s[idx:].count("]")
                    else:
                        combined.append(s)
                        depth += s.count("[") - s.count("]")
                        if depth <= 0:
                            break
                if combined:
                    try:
                        return json.loads("\n".join(combined))
                    except json.JSONDecodeError:
                        pass
                return None

            parsed = _find_json(raw)
            if parsed is None:
                logger.warning("[PPT] Parse failed. Raw: %s...", raw[:200])
                return []
            if isinstance(parsed, dict) and "slides" in parsed:
                parsed = parsed["slides"]
            if isinstance(parsed, list):
                logger.info("[PPT] Generated %d slides", len(parsed))
                return parsed
            return []
    except Exception as e:
        logger.warning("[PPT] Topic gen failed: %s", e)
        return []


def create_presentation_tool() -> ToolHandler:
    async def _execute(args: dict[str, Any], ctx: ToolContext | None) -> str:
        filename = args.get("filename", "presentation.pptx").strip()
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        title = args.get("title", "Presentation")

        # Read power model config at call time (after .env is loaded)
        power_model = os.environ.get("POWERPOINT_MODEL", "")
        power_key = os.environ.get("POWERPOINT_API_KEY", "") or os.environ.get("GROQ_API_KEY", "")
        power_url = os.environ.get("POWERPOINT_BASE_URL", "https://api.groq.com/openai/v1").rstrip("/")

        if power_model and power_key:
            logger.info("[PPT] Using power model %s for content generation", power_model)
            slides_data = await _generate_slides_from_topic(title, title, power_model, power_key, power_url)
            if not slides_data:
                return "ERROR: Failed to generate slides via power model."
        else:
            slides_data = args.get("slides", [])
            if not slides_data:
                return "ERROR: Provide 'slides' or set POWERPOINT_MODEL + POWERPOINT_API_KEY in .env."
        subtitle = args.get("subtitle", "")
        style_name = args.get("style", "corporate").lower()
        fmt = args.get("format", "16:9").lower()
        transition = args.get("transition", "").lower()

        st = STYLES.get(style_name, STYLES["corporate"])
        dims = CANVAS_FORMATS.get(fmt, CANVAS_FORMATS["16:9"])
        sw, sh = dims

        prs = PptxPresentation()
        prs.slide_width = sw
        prs.slide_height = sh

        # Title slide
        ts = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(ts, st.accent)
        bar = ts.shapes.add_shape(1, Inches(0), Inches(sh / 914400 - 0.15), sw, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = st.light if st.name != "Dark Tech" else RGBColor(0x00, 0x99, 0xCC)
        bar.line.fill.background()
        tf = _tb(ts, 1, 2.0, sw / 914400 - 2, 2.5).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(st.title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(22)
            p2.font.color.rgb = st.light
            p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = datetime.now().strftime("%B %d, %Y")
        p3.font.size = Pt(14)
        p3.font.color.rgb = st.light
        p3.alignment = PP_ALIGN.CENTER

        for sd in slides_data:
            stitle = sd.get("title", "Slide")
            content = sd.get("content", [])
            if isinstance(content, str):
                content = [content]
            layout = sd.get("layout")
            notes = sd.get("notes")
            img_src = sd.get("image", "") or None
            img_pos = sd.get("image_position", "right")
            cht = sd.get("chart")
            chd = sd.get("chart_data")
            await _build_slide(prs, stitle, content, layout, notes, st, sw, sh, img_src, img_pos, cht, chd)

        if transition in TRANSITIONS:
            try:
                for s in prs.slides:
                    if hasattr(s, 'slide_show_transition'):
                        s.slide_show_transition.transition_type = TRANSITIONS[transition]
            except Exception:
                pass

        SANDBOX.mkdir(parents=True, exist_ok=True)
        out = (SANDBOX / filename).resolve()
        prs.save(str(out))
        return f"Created: {filename} ({len(slides_data) + 1} slides, {st.name}, {fmt})"

    return ToolHandler(
        definition=ToolDefinition(
            function={
                "name": "create_presentation",
                "description": (
                    "Create a professional PowerPoint presentation."
                    " Provide 'title' and optionally 'slides', 'style',"
                    " 'transition'. Keep content short and structured."
                    " (If POWERPOINT_MODEL is set in .env, the tool"
                    " auto-generates better content using a stronger AI model.)"
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "Presentation title"},
                        "subtitle": {"type": "string", "description": "Optional subtitle"},
                        "filename": {"type": "string", "description": "Output filename (default: presentation.pptx)"},
                        "style": {
                            "type": "string",
                            "description": "Style: corporate, dark-tech, minimal, warm, nature, ocean, sunset, forest",
                        },
                        "format": {"type": "string", "description": "Format: 16:9 (default), 4:3, a4"},
                        "transition": {
                            "type": "string",
                            "description": "Transition: fade, push, wipe, cover, uncover, dissolve",
                        },
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "content": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": (
                                            "Slide content — each item is one line."
                                            " Keep bullets short (3-6 words) and parallel."
                                            " Auto-detects: bullets (default), code (```),"
                                            " quiz (?), two_column, timeline (time:desc),"
                                            " table (|col|col|), section, center_text."
                                            " Use # for sub-heading, > for quote."
                                        ),
                                    },
                                    "layout": {
                                        "type": "string",
                                        "description": (
                                            "Override: bullets, code, two_column,"
                                            " quiz, center_text, section,"
                                            " comparison, timeline, table"
                                        ),
                                    },
                                    "notes": {"type": "string"},
                                    "image": {"type": "string", "description": "URL or sandbox path for slide image"},
                                    "image_position": {
                                        "type": "string",
                                        "description": "right (default), left, full, bg (background)",
                                    },
                                    "chart": {"type": "string", "description": "Chart type: bar, column, line, pie"},
                                    "chart_data": {
                                        "type": "array",
                                        "items": {
                                            "type": "array",
                                            "items": {},
                                            "description": (
                                                "Chart rows as [label, value] pairs,"
                                                " e.g. [['Q1', 30], ['Q2', 50]]"
                                            ),
                                        },
                                    },
                                },
                                "required": ["title"],
                            },
                        },
                    },
                    "required": ["title"],
                },
            },
        ),
        execute=_execute,
    )
